# Copyright (c) 2026, one-fm and contributors
# Google Drive integration for Processa document storage (DMS).
#
# Credentials + config live on the "Processa Settings" singleton (one_bpmn),
# under its "Google Integration" section, resolved via
# integrations/google_common (shared with the Docs/Slides integrations and the
# connector layer):
#   google_drive_service_account_json  - service account JSON (Password field)
#
# (Destination folders are configured on the connector element / passed by the
# caller, not in settings.)
#
# Lookup order (google_common): Processa Settings -> AI Chat Settings (legacy
# fallback) -> frappe.conf -> <site>/private/files/gcp.json, so this also works
# from a bare bench console without the doctype populated.
#
# The backing GCP project must have the Drive API enabled, and the service
# account must be a member of the target Shared Drive (Content Manager or
# above) before any of this actually works — neither of those is something
# this module can do for itself.
#
# Called from Script Tasks (via one_bpmn's Server Script mechanism) — kept as
# plain reusable functions rather than one_bpmn-specific dispatchers, so the
# same functions serve any diagram's Script Task, not just one hardcoded
# BPMN service type.

import io
import json

import frappe

SCOPES = ["https://www.googleapis.com/auth/drive"]


class GoogleDriveConfigError(Exception):
	"""Raised for missing/invalid configuration — not a transient API failure."""


def _get_credentials(connector_id=None):
	# Credentials now come from the shared loader (Processa Settings → AI Chat
	# Settings → site_config → gcp.json). Re-raise as GoogleDriveConfigError so
	# existing Script Tasks that catch that type keep working unchanged.
	from one_bpmn.one_bpmn.integrations import google_common as _gc

	try:
		return _gc.get_credentials(SCOPES, connector_id=connector_id)
	except _gc.GoogleConfigError as e:
		raise GoogleDriveConfigError(str(e))


def _get_service(connector_id="google_drive"):
	"""Drive client, authenticated as the connector's own service account.

	Defaults to the google_drive connector so the Script Tasks that import these
	functions directly keep working without knowing about connectors at all.
	"""
	from googleapiclient.discovery import build

	return build(
		"drive", "v3", credentials=_get_credentials(connector_id), cache_discovery=False
	)


def create_file(
	folder_id: str,
	filename: str,
	content: str,
	target_mime_type: str = "application/vnd.google-apps.document",
	source_mime_type: str = "text/markdown",
) -> dict:
	"""
	Upload ``content`` into ``folder_id``, converting it into a native Google
	Doc (or whatever ``target_mime_type`` is) from ``source_mime_type``.

	Returns the created file's metadata: {"id", "name", "webViewLink"}.
	"""
	from googleapiclient.http import MediaInMemoryUpload

	service = _get_service()
	media = MediaInMemoryUpload(content.encode("utf-8"), mimetype=source_mime_type, resumable=False)
	file_metadata = {
		"name": filename,
		"parents": [folder_id],
		"mimeType": target_mime_type,
	}
	return (
		service.files()
		.create(
			body=file_metadata,
			media_body=media,
			fields="id,name,webViewLink",
			supportsAllDrives=True,
		)
		.execute()
	)


def update_file_content(file_id: str, content: str, source_mime_type: str = "text/markdown") -> dict:
	"""
	Replace an existing file's content — used to push finalized content into
	the empty placeholder file that create_file made earlier, once drafting/
	review is done. Returns the updated file's metadata.
	"""
	from googleapiclient.http import MediaInMemoryUpload

	service = _get_service()
	media = MediaInMemoryUpload(content.encode("utf-8"), mimetype=source_mime_type, resumable=False)
	return (
		service.files()
		.update(fileId=file_id, media_body=media, fields="id,name,webViewLink", supportsAllDrives=True)
		.execute()
	)


def set_permissions(file_id: str, grants: list) -> list:
	"""
	Apply a list of Drive permission grants to a file.

	Each grant is a dict per the Drive API Permissions resource, e.g.
	  {"type": "domain", "domain": "one-fm.com", "role": "reader"}
	  {"type": "user", "emailAddress": "x@one-fm.com", "role": "writer"}
	"""
	service = _get_service()
	results = []
	for grant in grants:
		results.append(
			service.permissions()
			.create(
				fileId=file_id,
				body=grant,
				fields="id",
				supportsAllDrives=True,
				sendNotificationEmail=False,
			)
			.execute()
		)
	return results


# Roles that express who *owns or structures* a file rather than who was granted
# access to read it. On a Shared Drive there is no "owner" role at all —
# ownership is the drive itself, surfaced as organizer/fileOrganizer — so
# skipping "owner" alone would have us trying to strip the drive's own
# management of the file. Drive refuses those deletions (403
# cannotDeletePermission), and it is right to: a file nobody organizes is
# orphaned.
STRUCTURAL_ROLES = ("owner", "organizer", "fileOrganizer")


def revoke_permissions(file_id: str, scope: str = "all", match: str = None) -> dict:
	"""Remove sharing grants from a file — the inverse of ``set_permissions``.

	This is what actually takes a document out of circulation. Marking a record
	inactive in Frappe hides it from Frappe; it does nothing about the
	domain-wide reader grant that ``set_permissions`` applied, so every employee
	holding the link still opens it. Revoking the grant is the part users see.

	``scope`` selects what to remove:
	  ``all``    — every access grant (the default: withdraw it from everyone)
	  ``domain`` — domain grants only, optionally just the one in ``match``
	  ``user``   — user grants only, optionally just the address in ``match``

	Returns ``{"removed": [...], "skipped": [...]}``. Two kinds of grant end up
	in ``skipped`` rather than failing the call:

	  * structural roles (see ``STRUCTURAL_ROLES``), which must not be removed;
	  * grants **inherited** from a parent folder, which Drive will not let you
	    delete on the item — they have to be changed where they are defined.

	Nothing is abandoned partway. Each deletion is attempted independently,
	because the alternative is what a single 403 produced here in practice: the
	domain grant stripped, the loop aborted, and a caller told the revocation
	failed when the important half had in fact succeeded. A partial withdrawal
	reported as a failure is far more dangerous than a complete one, so the
	caller gets the full picture instead.
	"""
	service = _get_service()
	listed = (
		service.permissions()
		.list(
			fileId=file_id,
			fields="permissions(id,type,role,emailAddress,domain,permissionDetails(inherited,inheritedFrom))",
			supportsAllDrives=True,
		)
		.execute()
	)

	removed = []
	skipped = []
	for perm in listed.get("permissions", []):
		if perm.get("role") in STRUCTURAL_ROLES:
			skipped.append(dict(perm, reason="structural role"))
			continue
		if any(d.get("inherited") for d in (perm.get("permissionDetails") or [])):
			skipped.append(dict(perm, reason="inherited from a parent folder"))
			continue
		if scope in ("domain", "user") and perm.get("type") != scope:
			continue
		if match:
			target = perm.get("domain") if perm.get("type") == "domain" else perm.get("emailAddress")
			if (target or "").lower() != match.lower():
				continue

		try:
			service.permissions().delete(
				fileId=file_id, permissionId=perm["id"], supportsAllDrives=True
			).execute()
			removed.append(perm)
		except Exception as e:
			skipped.append(dict(perm, reason=str(e)))

	return {"removed": removed, "skipped": skipped}


def list_files(folder_id: str, page_size: int = 20) -> list:
	"""List non-trashed files directly inside a Drive folder."""
	service = _get_service()
	resp = (
		service.files()
		.list(
			q=f"'{folder_id}' in parents and trashed = false",
			fields="files(id,name,mimeType,modifiedTime)",
			pageSize=page_size,
			supportsAllDrives=True,
			includeItemsFromAllDrives=True,
		)
		.execute()
	)
	return resp.get("files", [])


_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def _download_raw_bytes(service, file_id: str) -> bytes:
	from googleapiclient.http import MediaIoBaseDownload

	request = service.files().get_media(fileId=file_id)
	buf = io.BytesIO()
	downloader = MediaIoBaseDownload(buf, request)
	done = False
	while not done:
		_, done = downloader.next_chunk()
	return buf.getvalue()


def _extract_pptx_text(raw: bytes) -> str:
	from pptx import Presentation

	prs = Presentation(io.BytesIO(raw))
	blocks = []
	for i, slide in enumerate(prs.slides, start=1):
		lines = [f"## Slide {i}"]
		for shape in slide.shapes:
			if shape.has_text_frame:
				text = shape.text_frame.text.strip()
				if text:
					lines.append(text)
			if shape.has_table:
				for row in shape.table.rows:
					lines.append(" | ".join(cell.text.strip() for cell in row.cells))
		if len(lines) > 1:
			blocks.append("\n".join(lines))
	return "\n\n".join(blocks)


def _extract_docx_text(raw: bytes) -> str:
	from docx import Document as DocxDocument

	doc = DocxDocument(io.BytesIO(raw))
	blocks = [p.text for p in doc.paragraphs if p.text.strip()]
	for table in doc.tables:
		for row in table.rows:
			blocks.append(" | ".join(cell.text.strip() for cell in row.cells))
	return "\n".join(blocks)


def download_file_text(file_id: str, mime_type: str = None) -> str:
	"""
	Fetch a file's text content, converting to plain text regardless of
	source format:
	  - Native Google Docs/Slides   → Drive's own export-to-text API
	  - Raw .pptx (uploaded as-is)  → parsed slide-by-slide via python-pptx
	  - Raw .docx (uploaded as-is)  → parsed paragraph/table via python-docx
	  - Anything else (.md/.txt)    → best-effort UTF-8 decode

	mime_type is looked up from Drive automatically when not supplied —
	pass it explicitly only if the caller already has it on hand (e.g. from
	a prior list_files() call) to save the extra round trip.
	"""
	service = _get_service()
	if mime_type is None:
		mime_type = (
			service.files()
			.get(fileId=file_id, fields="mimeType", supportsAllDrives=True)
			.execute()
			.get("mimeType", "")
		)
	if mime_type in (
		"application/vnd.google-apps.document",
		"application/vnd.google-apps.presentation",
	):
		data = service.files().export(fileId=file_id, mimeType="text/plain").execute()
		return data.decode("utf-8") if isinstance(data, bytes) else data

	raw = _download_raw_bytes(service, file_id)

	# Match OOXML variants tolerantly: some editors (e.g. WPS Office) report
	# non-standard mimetypes like "application/wps-office.docx" — detect by the
	# real container instead of trusting the mimetype string.
	mt = (mime_type or "").lower()
	is_pptx = mt == _PPTX_MIME or mt.endswith("pptx") or "presentationml" in mt
	is_docx = mt == _DOCX_MIME or mt.endswith("docx") or "wordprocessingml" in mt

	# OOXML files are ZIP containers ("PK\x03\x04"); if the mimetype is unhelpful
	# but the bytes are a zip holding the tell-tale parts, sniff the type.
	if not (is_pptx or is_docx) and raw[:2] == b"PK":
		head = raw[:4000]
		if b"ppt/" in head:
			is_pptx = True
		elif b"word/" in head:
			is_docx = True

	if is_pptx:
		return _extract_pptx_text(raw)
	if is_docx:
		return _extract_docx_text(raw)

	return raw.decode("utf-8", errors="ignore")
